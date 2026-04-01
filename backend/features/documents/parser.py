import os
import re
import csv
import json
import zipfile
import tempfile
import shutil
try:
    import PyPDF2
except ImportError:
    PyPDF2 = None
try:
    from docx import Document
except ImportError:
    Document = None
try:
    import openpyxl
except ImportError:
    openpyxl = None
try:
    from pptx import Presentation
except ImportError:
    Presentation = None
try:
    from odf import opendocument
    from odf.table import Table
    from odf.text import P
    from odf.presentation import Slide
except ImportError:
    opendocument = None
    Table = None
    P = None
    Slide = None
try:
    import yaml
except ImportError:
    yaml = None
try:
    import email
    from email import policy
    from email.parser import BytesParser
except ImportError:
    email = None
try:
    from icalendar import Calendar
except ImportError:
    Calendar = None
try:
    import extract_msg
except ImportError:
    extract_msg = None
import markdown
try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None
import logging
import warnings

# PyPDF2 Warnings unterdrücken
if PyPDF2:
    warnings.filterwarnings("ignore", category=PyPDF2.errors.PdfReadWarning)

class DocumentParser:
    """Parser für verschiedene Dokumentformate"""
    
    def __init__(self):
        self.logger = logging.getLogger(__name__)
    
    def parse_file(self, file_path):
        """Parst eine Datei basierend auf ihrer Endung"""
        if not os.path.exists(file_path):
            self.logger.error(f"Datei nicht gefunden: {file_path}")
            return ""
        
        file_ext = os.path.splitext(file_path)[1].lower()
        
        try:
            if file_ext == '.pdf':
                return self.parse_pdf(file_path)
            elif file_ext in ['.docx', '.doc']:
                return self.parse_docx(file_path)
            elif file_ext in ['.xlsx', '.xls']:
                return self.parse_excel(file_path)
            elif file_ext in ['.pptx', '.ppt']:
                return self.parse_powerpoint(file_path)
            elif file_ext == '.odt':
                return self.parse_odt(file_path)
            elif file_ext == '.ods':
                return self.parse_ods(file_path)
            elif file_ext == '.odp':
                return self.parse_odp(file_path)
            elif file_ext in ['.md', '.markdown']:
                return self.parse_markdown(file_path)
            elif file_ext == '.csv':
                return self.parse_csv(file_path)
            elif file_ext == '.json':
                return self.parse_json(file_path)
            elif file_ext in ['.yaml', '.yml']:
                return self.parse_yaml(file_path)
            elif file_ext == '.xml':
                return self.parse_xml(file_path)
            elif file_ext == '.eml':
                return self.parse_eml(file_path)
            elif file_ext == '.msg':
                return self.parse_msg(file_path)
            elif file_ext == '.ics':
                return self.parse_ics(file_path)
            elif file_ext == '.zip':
                return self.parse_zip(file_path)
            elif file_ext == '.txt':
                return self.parse_text(file_path)
            elif file_ext in ['.html', '.htm']:
                return self.parse_html(file_path)
            else:
                self.logger.warning(f"Unsupported format: {file_ext}")
                return ""
        except Exception as e:
            self.logger.error(f"Fehler beim Parsen von {file_path}: {str(e)}")
            return ""
    
    def parse_pdf(self, file_path):
        """Parst PDF-Dateien mit robuster Fehlerbehandlung"""
        text = ""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file, strict=False)  # strict=False für mehr Toleranz
                
                for page_num in range(len(pdf_reader.pages)):
                    try:
                        page = pdf_reader.pages[page_num]
                        page_text = page.extract_text()
                        if page_text and page_text.strip():  # Nur hinzufügen wenn Text vorhanden
                            text += page_text + "\n"
                    except Exception as page_error:
                        self.logger.warning(f"Error reading page {page_num + 1} in {file_path}: {str(page_error)}")
                        continue
                        
        except Exception as e:
            self.logger.error(f"PDF parsing error for {file_path}: {str(e)}")
        
        return self.clean_text(text)
    
    def parse_docx(self, file_path):
        """Parst Word-Dokumente"""
        text = ""
        try:
            doc = Document(file_path)
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            
            # Tabellen extrahieren
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        row_text.append(cell.text.strip())
                    text += " | ".join(row_text) + "\n"
        except Exception as e:
            self.logger.error(f"DOCX parsing error: {str(e)}")
        
        return self.clean_text(text)
    
    def parse_excel(self, file_path):
        """Parst Excel-Dateien"""
        text = ""
        try:
            workbook = openpyxl.load_workbook(file_path, read_only=True)
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"\n=== Sheet: {sheet_name} ===\n"
                
                for row in sheet.iter_rows(values_only=True):
                    row_data = [str(cell) if cell is not None else "" for cell in row]
                    if any(row_data):  # Leere Zeilen überspringen
                        text += " | ".join(row_data) + "\n"
        except Exception as e:
            self.logger.error(f"Excel parsing error: {str(e)}")
        
        return self.clean_text(text)
    
    def parse_powerpoint(self, file_path):
        """Parst PowerPoint-Präsentationen"""
        text = ""
        try:
            prs = Presentation(file_path)
            for slide_num, slide in enumerate(prs.slides):
                text += f"\n=== Slide {slide_num + 1} ===\n"
                
                # Text aus Shapes extrahieren
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
                    
                    # Tabellen in Slides
                    if shape.has_table:
                        table = shape.table
                        for row_idx, row in enumerate(table.rows):
                            row_text = []
                            for cell in row.cells:
                                row_text.append(cell.text.strip())
                            text += " | ".join(row_text) + "\n"
        except Exception as e:
            self.logger.error(f"PowerPoint parsing error: {str(e)}")

        return self.clean_text(text)

    def parse_odt(self, file_path):
        """Parst OpenDocument Text-Dateien"""
        text = ""
        try:
            doc = opendocument.load(file_path)
            paragraphs = doc.getElementsByType(P)
            for p in paragraphs:
                for node in p.childNodes:
                    if hasattr(node, 'data'):
                        text += node.data + "\n"
        except Exception as e:
            self.logger.error(f"ODT parsing error: {str(e)}")
        return self.clean_text(text)

    def parse_ods(self, file_path):
        """Parst OpenDocument Spreadsheet-Dateien"""
        text = ""
        try:
            doc = opendocument.load(file_path)
            tables = doc.getElementsByType(Table)
            for table in tables:
                text += "\n=== Tabelle ===\n"
                for row in table.getElementsByType(opendocument.table.TableRow):
                    row_data = []
                    for cell in row.getElementsByType(opendocument.table.TableCell):
                        cell_text = ""
                        for p in cell.getElementsByType(P):
                            for node in p.childNodes:
                                if hasattr(node, 'data'):
                                    cell_text += node.data
                        row_data.append(cell_text)
                    if any(row_data):
                        text += " | ".join(row_data) + "\n"
        except Exception as e:
            self.logger.error(f"ODS parsing error: {str(e)}")
        return self.clean_text(text)

    def parse_odp(self, file_path):
        """Parst OpenDocument Presentation-Dateien"""
        text = ""
        try:
            doc = opendocument.load(file_path)
            slides = doc.getElementsByType(opendocument.presentation.Page)
            for slide_num, slide in enumerate(slides, 1):
                text += f"\n=== Folie {slide_num} ===\n"
                # Text aus allen Absätzen extrahieren
                for p in slide.getElementsByType(P):
                    for node in p.childNodes:
                        if hasattr(node, 'data'):
                            text += node.data + "\n"
        except Exception as e:
            self.logger.error(f"ODP parsing error: {str(e)}")
        return self.clean_text(text)
    
    def parse_markdown(self, file_path):
        """Parst Markdown-Dateien"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                md_content = file.read()
            
            # Markdown zu HTML konvertieren und dann Text extrahieren
            html = markdown.markdown(md_content)
            soup = BeautifulSoup(html, 'html.parser')
            text = soup.get_text()
            
            return self.clean_text(text)
        except Exception as e:
            self.logger.error(f"Markdown parsing error: {str(e)}")
            return ""
    
    def parse_text(self, file_path):
        """Parst reine Textdateien"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return self.clean_text(file.read())
        except UnicodeDecodeError:
            # Versuche mit anderer Kodierung
            try:
                with open(file_path, 'r', encoding='latin-1') as file:
                    return self.clean_text(file.read())
            except Exception as e:
                self.logger.error(f"Text parsing error: {str(e)}")
                return ""
        except Exception as e:
            self.logger.error(f"Text parsing error: {str(e)}")
            return ""
    
    def parse_html(self, file_path):
        """Parst HTML-Dateien"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                html_content = file.read()
            
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Scripts und Styles entfernen
            for script in soup(["script", "style"]):
                script.decompose()
            
            text = soup.get_text()
            return self.clean_text(text)
        except Exception as e:
            self.logger.error(f"HTML parsing error: {str(e)}")
            return ""

    def parse_csv(self, file_path):
        """Parst CSV-Dateien"""
        text = ""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                reader = csv.reader(f)
                for row in reader:
                    if any(cell.strip() for cell in row):
                        text += " | ".join(row) + "\n"
        except Exception as e:
            self.logger.error(f"CSV parsing error: {str(e)}")
        return self.clean_text(text)

    def parse_json(self, file_path):
        """Parst JSON-Dateien"""
        text = ""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                data = json.load(f)
            # Rekursive Funktion zum Extrahieren von Strings
            def extract_strings(obj, path=""):
                result = []
                if isinstance(obj, dict):
                    for k, v in obj.items():
                        result.extend(extract_strings(v, f"{path}.{k}" if path else k))
                elif isinstance(obj, list):
                    for i, item in enumerate(obj):
                        result.extend(extract_strings(item, f"{path}[{i}]"))
                elif isinstance(obj, str):
                    result.append(f"{path}: {obj}")
                elif isinstance(obj, (int, float, bool)):
                    result.append(f"{path}: {obj}")
                return result
            strings = extract_strings(data)
            text = "\n".join(strings)
        except Exception as e:
            self.logger.error(f"JSON parsing error: {str(e)}")
        return self.clean_text(text)

    def parse_yaml(self, file_path):
        """Parst YAML-Dateien"""
        if not yaml:
            self.logger.warning("yaml module not installed")
            return ""
        text = ""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                data = yaml.safe_load(f)
            # Ähnlich wie JSON: Strings extrahieren
            def extract_strings(obj, path=""):
                result = []
                if isinstance(obj, dict):
                    for k, v in obj.items():
                        result.extend(extract_strings(v, f"{path}.{k}" if path else k))
                elif isinstance(obj, list):
                    for i, item in enumerate(obj):
                        result.extend(extract_strings(item, f"{path}[{i}]"))
                elif isinstance(obj, str):
                    result.append(f"{path}: {obj}")
                elif isinstance(obj, (int, float, bool)):
                    result.append(f"{path}: {obj}")
                return result
            if data:
                strings = extract_strings(data)
                text = "\n".join(strings)
        except Exception as e:
            self.logger.error(f"YAML parsing error: {str(e)}")
        return self.clean_text(text)

    def parse_xml(self, file_path):
        """Parst XML-Dateien"""
        text = ""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            if BeautifulSoup:
                soup = BeautifulSoup(content, 'xml')
                text = soup.get_text()
            else:
                # Fallback: lxml oder einfaches Strippen
                import xml.etree.ElementTree as ET
                root = ET.fromstring(content)
                def get_text(elem):
                    result = []
                    if elem.text and elem.text.strip():
                        result.append(elem.text.strip())
                    for child in elem:
                        result.extend(get_text(child))
                        if child.tail and child.tail.strip():
                            result.append(child.tail.strip())
                    return result
                text = "\n".join(get_text(root))
        except Exception as e:
            self.logger.error(f"XML parsing error: {str(e)}")
        return self.clean_text(text)

    def parse_eml(self, file_path):
        """Parst EML-Dateien (E-Mail)"""
        if not email:
            self.logger.warning("email module not available")
            return ""
        text = ""
        try:
            with open(file_path, 'rb') as f:
                msg = BytesParser(policy=policy.default).parse(f)
            # Header
            text += f"Von: {msg.get('From', '')}\n"
            text += f"An: {msg.get('To', '')}\n"
            text += f"Betreff: {msg.get('Subject', '')}\n"
            text += f"Datum: {msg.get('Date', '')}\n\n"
            # Body
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == 'text/plain':
                        try:
                            text += part.get_content()
                        except:
                            pass
            else:
                if msg.get_content_type() == 'text/plain':
                    text += msg.get_content()
        except Exception as e:
            self.logger.error(f"EML parsing error: {str(e)}")
        return self.clean_text(text)

    def parse_msg(self, file_path):
        """Parst MSG-Dateien (Outlook)"""
        if not extract_msg:
            self.logger.warning("extract_msg module not installed")
            return ""
        text = ""
        try:
            msg = extract_msg.Message(file_path)
            text += f"Von: {msg.sender}\n"
            text += f"An: {msg.to}\n"
            text += f"Betreff: {msg.subject}\n"
            text += f"Datum: {msg.date}\n\n"
            text += msg.body
        except Exception as e:
            self.logger.error(f"MSG parsing error: {str(e)}")
        return self.clean_text(text)

    def parse_ics(self, file_path):
        """Parst ICS-Dateien (iCalendar)"""
        if not Calendar:
            self.logger.warning("icalendar module not installed")
            return ""
        text = ""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                cal = Calendar.from_ical(f.read())
            for component in cal.walk():
                if component.name == "VEVENT":
                    text += "\n=== Event ===\n"
                    summary = component.get('summary')
                    if summary:
                        text += f"Titel: {summary}\n"
                    start = component.get('dtstart')
                    if start:
                        text += f"Start: {start.dt}\n"
                    end = component.get('dtend')
                    if end:
                        text += f"Ende: {end.dt}\n"
                    desc = component.get('description')
                    if desc:
                        text += f"Beschreibung: {desc}\n"
                    loc = component.get('location')
                    if loc:
                        text += f"Ort: {loc}\n"
        except Exception as e:
            self.logger.error(f"ICS parsing error: {str(e)}")
        return self.clean_text(text)

    def parse_zip(self, file_path):
        """Parst ZIP-Dateien (rekursiv alle enthaltenen Dateien)"""
        text = ""
        temp_dir = None
        try:
            temp_dir = tempfile.mkdtemp()
            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                for member in zip_ref.infolist():
                    name = member.filename
                    text += f"\n=== Datei in ZIP: {name} ===\n"

                    # Skip directories and unsafe absolute/parent paths.
                    if member.is_dir() or os.path.isabs(name) or '..' in name.split('/'):
                        self.logger.warning(f"Skipping unsafe ZIP member: {name}")
                        continue

                    extracted_path = os.path.realpath(os.path.join(temp_dir, name))
                    temp_root = os.path.realpath(temp_dir)
                    if not extracted_path.startswith(temp_root + os.sep):
                        self.logger.warning(f"Skipping ZIP traversal attempt: {name}")
                        continue

                    try:
                        zip_ref.extract(member, temp_dir)
                        file_text = self.parse_file(extracted_path)
                        if file_text:
                            text += file_text + "\n"
                    except Exception as e:
                        self.logger.warning(f"Error parsing {name} in ZIP: {str(e)}")
        except Exception as e:
            self.logger.error(f"ZIP parsing error: {str(e)}")
        finally:
            if temp_dir and os.path.exists(temp_dir):
                shutil.rmtree(temp_dir)
        return self.clean_text(text)
    
    def clean_text(self, text):
        """Bereinigt den extrahierten Text"""
        if not text:
            return ""
        
        # Mehrfache Leerzeichen entfernen
        text = re.sub(r'\s+', ' ', text)
        
        # Zeilenumbrüche normalisieren
        text = re.sub(r'\n\s*\n', '\n\n', text)
        
        # Sonderzeichen bereinigen
        text = text.strip()
        
        # Debug-Info hinzufügen
        if len(text) < 50:
            self.logger.debug(f"Short text extracted: '{text[:100]}...' (length: {len(text)})")
        
        return text
