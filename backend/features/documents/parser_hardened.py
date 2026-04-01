"""
Hardened Document Parser with security fixes.

Security improvements:
- ZIP path traversal prevention (CWE-22)
- File size limits
- Safe temporary file handling
- Input validation
- Resource exhaustion prevention
"""

import os
import re
import csv
import json
import zipfile
import tempfile
import shutil
import logging
from typing import Optional

# Optional dependencies with graceful degradation
try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from odf import opendocument
    from odf.table import Table
    from odf.text import P
    from odf.presentation import Slide
except ImportError:
    opendocument = None

try:
    import yaml
except ImportError:
    yaml = None

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

import markdown
import warnings

if PyPDF2:
    warnings.filterwarnings("ignore", category=PyPDF2.errors.PdfReadWarning)


class DocumentParserError(Exception):
    """Base exception for document parsing errors."""
    pass


class FileSizeError(DocumentParserError):
    """File exceeds size limits."""
    pass


class PathTraversalError(DocumentParserError):
    """Path traversal attempt detected."""
    pass


class DocumentParser:
    """
    Hardened parser for multiple document formats.

    Security features:
    - Path traversal prevention in ZIP files
    - File size limits (configurable)
    - Safe temporary file handling
    - Resource exhaustion prevention
    """

    # Security constraints
    MAX_FILE_SIZE = 50 * 1024 * 1024  # 50 MB
    MAX_EXTRACTED_SIZE = 200 * 1024 * 1024  # 200 MB total
    MAX_ZIP_MEMBERS = 1000
    TEMP_DIR_PREFIX = "mynd_parser_"

    def __init__(self, max_file_size: int = None):
        """
        Initialize parser.

        Args:
            max_file_size: Max individual file size (bytes)
        """
        self.logger = logging.getLogger(__name__)
        self.max_file_size = max_file_size or self.MAX_FILE_SIZE

    def parse_file(self, file_path: str) -> str:
        """
        Parse file to text.

        Args:
            file_path: Path to file

        Returns:
            Extracted text

        Raises:
            DocumentParserError: If parsing fails
        """
        if not os.path.exists(file_path):
            raise DocumentParserError(f"File not found: {file_path}")

        # Check file size
        try:
            file_size = os.path.getsize(file_path)
            if file_size > self.max_file_size:
                raise FileSizeError(
                    f"File too large: {file_size} > {self.max_file_size}"
                )
        except OSError as e:
            raise DocumentParserError(f"Cannot stat file: {e}")

        file_ext = os.path.splitext(file_path)[1].lower()

        try:
            if file_ext == ".pdf":
                return self.parse_pdf(file_path)
            elif file_ext in (".docx", ".doc"):
                return self.parse_docx(file_path)
            elif file_ext in (".xlsx", ".xls"):
                return self.parse_excel(file_path)
            elif file_ext in (".pptx", ".ppt"):
                return self.parse_powerpoint(file_path)
            elif file_ext in (".md", ".markdown"):
                return self.parse_markdown(file_path)
            elif file_ext == ".csv":
                return self.parse_csv(file_path)
            elif file_ext == ".json":
                return self.parse_json(file_path)
            elif file_ext in (".html", ".htm"):
                return self.parse_html(file_path)
            elif file_ext == ".zip":
                return self.parse_zip_secure(file_path)  # Use hardened version
            elif file_ext == ".txt":
                return self.parse_text(file_path)
            else:
                self.logger.warning(f"Unsupported format: {file_ext}")
                return ""
        except DocumentParserError:
            raise
        except Exception as e:
            self.logger.error(f"Parse error for {file_path}: {type(e).__name__}")
            raise DocumentParserError(f"Failed to parse {file_path}: {e}")

    def parse_pdf(self, file_path: str) -> str:
        """Parse PDF file."""
        if not PyPDF2:
            raise DocumentParserError("PyPDF2 not installed")

        text = ""
        try:
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f, strict=False)
                for page_num, page in enumerate(reader.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text and page_text.strip():
                            text += page_text + "\n"
                    except Exception as e:
                        self.logger.warning(f"Error reading page {page_num}: {e}")
                        continue
        except Exception as e:
            raise DocumentParserError(f"PDF parse error: {e}")

        return self.clean_text(text)

    def parse_docx(self, file_path: str) -> str:
        """Parse DOCX file."""
        if not Document:
            raise DocumentParserError("python-docx not installed")

        text = ""
        try:
            doc = Document(file_path)
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"

            for table in doc.tables:
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    text += " | ".join(row_text) + "\n"
        except Exception as e:
            raise DocumentParserError(f"DOCX parse error: {e}")

        return self.clean_text(text)

    def parse_excel(self, file_path: str) -> str:
        """Parse Excel file."""
        if not openpyxl:
            raise DocumentParserError("openpyxl not installed")

        text = ""
        try:
            workbook = openpyxl.load_workbook(file_path, read_only=True)
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"\n=== Sheet: {sheet_name} ===\n"

                for row in sheet.iter_rows(values_only=True):
                    row_data = [str(cell) if cell else "" for cell in row]
                    if any(row_data):
                        text += " | ".join(row_data) + "\n"
        except Exception as e:
            raise DocumentParserError(f"Excel parse error: {e}")

        return self.clean_text(text)

    def parse_powerpoint(self, file_path: str) -> str:
        """Parse PowerPoint file."""
        if not Presentation:
            raise DocumentParserError("python-pptx not installed")

        text = ""
        try:
            prs = Presentation(file_path)
            for slide_num, slide in enumerate(prs.slides):
                text += f"\n=== Slide {slide_num + 1} ===\n"

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"

                    if hasattr(shape, "has_table") and shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = [cell.text for cell in row.cells]
                            text += " | ".join(row_text) + "\n"
        except Exception as e:
            raise DocumentParserError(f"PowerPoint parse error: {e}")

        return self.clean_text(text)

    def parse_markdown(self, file_path: str) -> str:
        """Parse Markdown file."""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                md_content = f.read()

            # Convert MD to HTML then extract text
            html = markdown.markdown(md_content)
            if BeautifulSoup:
                soup = BeautifulSoup(html, "html.parser")
                text = soup.get_text()
            else:
                # Fallback: strip HTML tags
                text = re.sub(r"<[^>]+>", "", html)

            return self.clean_text(text)
        except Exception as e:
            raise DocumentParserError(f"Markdown parse error: {e}")

    def parse_csv(self, file_path: str) -> str:
        """Parse CSV file."""
        text = ""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                reader = csv.reader(f)
                for row in reader:
                    text += " | ".join(str(cell) for cell in row) + "\n"
        except Exception as e:
            raise DocumentParserError(f"CSV parse error: {e}")

        return self.clean_text(text)

    def parse_json(self, file_path: str) -> str:
        """Parse JSON file."""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                data = json.load(f)

            # Convert to readable format
            text = json.dumps(data, indent=2, ensure_ascii=False)
            return self.clean_text(text)
        except Exception as e:
            raise DocumentParserError(f"JSON parse error: {e}")

    def parse_html(self, file_path: str) -> str:
        """Parse HTML file."""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                html_content = f.read()

            if BeautifulSoup:
                soup = BeautifulSoup(html_content, "html.parser")
                # Remove script and style elements
                for script in soup(["script", "style"]):
                    script.decompose()
                text = soup.get_text()
            else:
                # Fallback: strip tags
                text = re.sub(r"<[^>]+>", "", html_content)

            return self.clean_text(text)
        except Exception as e:
            raise DocumentParserError(f"HTML parse error: {e}")

    def parse_text(self, file_path: str) -> str:
        """Parse plain text file."""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()
            return self.clean_text(text)
        except Exception as e:
            raise DocumentParserError(f"Text parse error: {e}")

    def parse_zip_secure(self, file_path: str) -> str:
        """
        Parse ZIP file with path traversal prevention (CWE-22).

        Security:
        - Validates all extracted paths
        - Limits number of members
        - Enforces size limits
        - Safe temporary directory
        """
        text = ""
        temp_dir = None

        try:
            # Use secure temporary directory
            temp_dir = tempfile.mkdtemp(prefix=self.TEMP_DIR_PREFIX)
            os.chmod(temp_dir, 0o700)  # Only owner can read

            extracted_size = 0

            with zipfile.ZipFile(file_path, "r") as zip_ref:
                members = zip_ref.infolist()

                # Limit number of members
                if len(members) > self.MAX_ZIP_MEMBERS:
                    raise FileSizeError(
                        f"ZIP contains too many members: {len(members)} > {self.MAX_ZIP_MEMBERS}"
                    )

                for member in members:
                    # Prevent directory traversal
                    try:
                        self._validate_zip_member(member, temp_dir)
                    except PathTraversalError as e:
                        self.logger.warning(f"Skipping unsafe member: {e}")
                        continue

                    # Check size
                    if member.file_size > self.max_file_size:
                        self.logger.warning(
                            f"Skipping large member: {member.filename}"
                        )
                        continue

                    extracted_size += member.file_size
                    if extracted_size > self.MAX_EXTRACTED_SIZE:
                        raise FileSizeError(
                            f"Extracted size exceeds limit: {extracted_size}"
                        )

                    # Extract safely
                    try:
                        zip_ref.extract(member, temp_dir)
                        extracted_path = os.path.join(temp_dir, member.filename)

                        if os.path.isfile(extracted_path):
                            text += f"\n=== {member.filename} ===\n"
                            try:
                                content = self.parse_file(extracted_path)
                                if content:
                                    text += content + "\n"
                            except DocumentParserError:
                                pass
                    except Exception as e:
                        self.logger.warning(
                            f"Failed to extract {member.filename}: {e}"
                        )

            return self.clean_text(text)

        except (FileSizeError, PathTraversalError):
            raise
        except Exception as e:
            raise DocumentParserError(f"ZIP parse error: {e}")
        finally:
            if temp_dir and os.path.exists(temp_dir):
                shutil.rmtree(temp_dir, ignore_errors=True)

    def _validate_zip_member(self, member, base_dir: str) -> None:
        """
        Validate ZIP member path (CWE-22 prevention).

        Args:
            member: ZipInfo member
            base_dir: Base extraction directory

        Raises:
            PathTraversalError: If path escape detected
        """
        name = member.filename

        # Reject absolute paths
        if os.path.isabs(name):
            raise PathTraversalError(f"Absolute path: {name}")

        # Reject parent directory references
        if ".." in name or name.startswith("/"):
            raise PathTraversalError(f"Parent directory reference: {name}")

        # Get real paths
        extract_path = os.path.realpath(os.path.join(base_dir, name))
        base_real = os.path.realpath(base_dir)

        # Ensure extraction stays within base_dir
        common = os.path.commonpath([base_real, extract_path])
        if common != base_real:
            raise PathTraversalError(
                f"Path traversal attempt: {name} -> {extract_path}"
            )

    def clean_text(self, text: str) -> str:
        """Clean and normalize extracted text."""
        if not text:
            return ""

        # Remove control characters
        text = "".join(c for c in text if ord(c) >= 32 or c in "\n\r\t")

        # Normalize whitespace
        text = re.sub(r"\s+", " ", text)
        text = re.sub(r"\n\s*\n", "\n\n", text)

        return text.strip()
