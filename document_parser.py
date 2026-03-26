import os
import re
try:
    import PyPDF2
except ImportError:
    PyPDF2 = None
try:
    from docx import Document
except ImportError:
    Document = None
try:
    import openpyxl
except ImportError:
    openpyxl = None
try:
    from pptx import Presentation
except ImportError:
    Presentation = None
import markdown
try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None
import logging
import warnings

# PyPDF2 Warnings unterdrücken
if PyPDF2:
    warnings.filterwarnings("ignore", category=PyPDF2.errors.PdfReadWarning)

class DocumentParser:
    """Parser für verschiedene Dokumentformate"""
    
    def __init__(self):
        self.logger = logging.getLogger(__name__)
    
    def parse_file(self, file_path):
        """Parst eine Datei basierend auf ihrer Endung"""
        if not os.path.exists(file_path):
            self.logger.error(f"Datei nicht gefunden: {file_path}")
            return ""
        
        file_ext = os.path.splitext(file_path)[1].lower()
        
        try:
            if file_ext == '.pdf':
                return self.parse_pdf(file_path)
            elif file_ext in ['.docx', '.doc']:
                return self.parse_docx(file_path)
            elif file_ext in ['.xlsx', '.xls']:
                return self.parse_excel(file_path)
            elif file_ext in ['.pptx', '.ppt']:
                return self.parse_powerpoint(file_path)
            elif file_ext in ['.md', '.markdown']:
                return self.parse_markdown(file_path)
            elif file_ext == '.txt':
                return self.parse_text(file_path)
            elif file_ext in ['.html', '.htm']:
                return self.parse_html(file_path)
            else:
                self.logger.warning(f"Unsupported format: {file_ext}")
                return ""
        except Exception as e:
            self.logger.error(f"Fehler beim Parsen von {file_path}: {str(e)}")
            return ""
    
    def parse_pdf(self, file_path):
        """Parst PDF-Dateien mit robuster Fehlerbehandlung"""
        text = ""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file, strict=False)  # strict=False für mehr Toleranz
                
                for page_num in range(len(pdf_reader.pages)):
                    try:
                        page = pdf_reader.pages[page_num]
                        page_text = page.extract_text()
                        if page_text and page_text.strip():  # Nur hinzufügen wenn Text vorhanden
                            text += page_text + "\n"
                    except Exception as page_error:
                        self.logger.warning(f"Error reading page {page_num + 1} in {file_path}: {str(page_error)}")
                        continue
                        
        except Exception as e:
            self.logger.error(f"PDF parsing error for {file_path}: {str(e)}")
        
        return self.clean_text(text)
    
    def parse_docx(self, file_path):
        """Parst Word-Dokumente"""
        text = ""
        try:
            doc = Document(file_path)
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            
            # Tabellen extrahieren
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        row_text.append(cell.text.strip())
                    text += " | ".join(row_text) + "\n"
        except Exception as e:
            self.logger.error(f"DOCX parsing error: {str(e)}")
        
        return self.clean_text(text)
    
    def parse_excel(self, file_path):
        """Parst Excel-Dateien"""
        text = ""
        try:
            workbook = openpyxl.load_workbook(file_path, read_only=True)
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"\n=== Sheet: {sheet_name} ===\n"
                
                for row in sheet.iter_rows(values_only=True):
                    row_data = [str(cell) if cell is not None else "" for cell in row]
                    if any(row_data):  # Leere Zeilen überspringen
                        text += " | ".join(row_data) + "\n"
        except Exception as e:
            self.logger.error(f"Excel parsing error: {str(e)}")
        
        return self.clean_text(text)
    
    def parse_powerpoint(self, file_path):
        """Parst PowerPoint-Präsentationen"""
        text = ""
        try:
            prs = Presentation(file_path)
            for slide_num, slide in enumerate(prs.slides):
                text += f"\n=== Slide {slide_num + 1} ===\n"
                
                # Text aus Shapes extrahieren
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
                    
                    # Tabellen in Slides
                    if shape.has_table:
                        table = shape.table
                        for row_idx, row in enumerate(table.rows):
                            row_text = []
                            for cell in row.cells:
                                row_text.append(cell.text.strip())
                            text += " | ".join(row_text) + "\n"
        except Exception as e:
            self.logger.error(f"PowerPoint parsing error: {str(e)}")
        
        return self.clean_text(text)
    
    def parse_markdown(self, file_path):
        """Parst Markdown-Dateien"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                md_content = file.read()
            
            # Markdown zu HTML konvertieren und dann Text extrahieren
            html = markdown.markdown(md_content)
            soup = BeautifulSoup(html, 'html.parser')
            text = soup.get_text()
            
            return self.clean_text(text)
        except Exception as e:
            self.logger.error(f"Markdown parsing error: {str(e)}")
            return ""
    
    def parse_text(self, file_path):
        """Parst reine Textdateien"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return self.clean_text(file.read())
        except UnicodeDecodeError:
            # Versuche mit anderer Kodierung
            try:
                with open(file_path, 'r', encoding='latin-1') as file:
                    return self.clean_text(file.read())
            except Exception as e:
                self.logger.error(f"Text parsing error: {str(e)}")
                return ""
        except Exception as e:
            self.logger.error(f"Text parsing error: {str(e)}")
            return ""
    
    def parse_html(self, file_path):
        """Parst HTML-Dateien"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                html_content = file.read()
            
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Scripts und Styles entfernen
            for script in soup(["script", "style"]):
                script.decompose()
            
            text = soup.get_text()
            return self.clean_text(text)
        except Exception as e:
            self.logger.error(f"HTML parsing error: {str(e)}")
            return ""
    
    def clean_text(self, text):
        """Bereinigt den extrahierten Text"""
        if not text:
            return ""
        
        # Mehrfache Leerzeichen entfernen
        text = re.sub(r'\s+', ' ', text)
        
        # Zeilenumbrüche normalisieren
        text = re.sub(r'\n\s*\n', '\n\n', text)
        
        # Sonderzeichen bereinigen
        text = text.strip()
        
        # Debug-Info hinzufügen
        if len(text) < 50:
            self.logger.debug(f"Short text extracted: '{text[:100]}...' (length: {len(text)})")
        
        return text
